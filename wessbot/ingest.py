"""Document extraction, chunking, and Chroma ingestion utilities."""
from __future__ import annotations

from dataclasses import dataclass
import argparse
import hashlib
import json
from pathlib import Path
from typing import Iterable

import chromadb
from openai import OpenAI

from .config import BASE_DIR, CHROMA_DIR, DEFAULT_EMBEDDING_MODEL
from .products import PRODUCTS


@dataclass(frozen=True)
class SourceGroup:
    product: str
    collection: str
    paths: tuple[Path, ...]


def default_source_groups(base_dir: Path = BASE_DIR) -> list[SourceGroup]:
    parent = base_dir.parent
    return [
        SourceGroup("ENV200", PRODUCTS["ENV200"].collection, (base_dir / "docs", parent / "농도계 (ENV200)")),
        SourceGroup("ENV130", PRODUCTS["ENV130"].collection, (base_dir / "docs_interface", parent / "계면계 (ENV130)")),
        SourceGroup("ENV120", PRODUCTS["ENV120"].collection, (base_dir / "docs_env120", parent / "계면계 (ENV120)")),
    ]


def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def extract_text_from_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_text_from_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    texts: list[str] = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            texts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                texts.append(" | ".join(cells))
    return "\n".join(texts)


def extract_text_from_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            texts.append(f"[Slide {slide_no}]\n" + "\n".join(slide_texts))
    return "\n".join(texts)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".md":
        return extract_text_from_md(path)
    if suffix == ".docx":
        return extract_text_from_docx(path)
    if suffix == ".pptx":
        return extract_text_from_pptx(path)
    raise ValueError(f"Unsupported file type: {path}")


def chunk_text(text: str, chunk_size: int = 900, overlap: int = 120) -> list[str]:
    """Chunk text by paragraphs with character overlap.

    The larger chunk size gives the support model enough local context while still
    keeping retrieval precise for menu/procedure questions.
    """
    paragraphs = [p.strip() for p in text.replace("\r\n", "\n").split("\n") if p.strip()]
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0
    for para in paragraphs:
        para_len = len(para)
        if current and current_len + para_len + 1 > chunk_size:
            chunks.append("\n".join(current).strip())
            overlap_lines: list[str] = []
            overlap_len = 0
            for prev in reversed(current):
                if overlap_len + len(prev) > overlap:
                    break
                overlap_lines.insert(0, prev)
                overlap_len += len(prev)
            current = overlap_lines
            current_len = overlap_len
        current.append(para)
        current_len += para_len + 1
    if current:
        chunks.append("\n".join(current).strip())
    return [c for c in chunks if c]


def iter_source_files(paths: Iterable[Path]) -> list[Path]:
    files: list[Path] = []
    for root in paths:
        if not root.exists():
            continue
        for path in sorted(root.rglob("*")):
            if path.is_file() and path.suffix.lower() in {".md", ".docx", ".pptx"} and not path.name.startswith("~$"):
                files.append(path)
    return files


def embed_text(client: OpenAI, text: str, model: str) -> list[float]:
    return client.embeddings.create(model=model, input=text).data[0].embedding


def rebuild_database(
    *,
    chroma_dir: Path = CHROMA_DIR,
    base_dir: Path = BASE_DIR,
    embedding_model: str = DEFAULT_EMBEDDING_MODEL,
    chunk_size: int = 900,
    overlap: int = 120,
    dry_run: bool = False,
    allow_partial: bool = False,
) -> dict:
    client = None if dry_run else OpenAI()
    chroma_client = None if dry_run else chromadb.PersistentClient(path=str(chroma_dir))
    report: dict = {"chroma_dir": str(chroma_dir), "embedding_model": embedding_model, "products": {}}

    for group in default_source_groups(base_dir):
        files = iter_source_files(group.paths)
        missing_paths = [str(p) for p in group.paths if not p.exists()]
        if missing_paths and not dry_run and not allow_partial:
            raise RuntimeError(
                f"Refusing to rebuild {group.collection}: missing source paths {missing_paths}. "
                "Run with --dry-run first, restore the source document folders, or pass --allow-partial if this is intentional."
            )
        product_report = {"collection": group.collection, "files": [], "total_chunks": 0, "missing_paths": missing_paths}
        if not dry_run:
            try:
                chroma_client.delete_collection(group.collection)
            except Exception:
                pass
            collection = chroma_client.create_collection(name=group.collection, metadata={"hnsw:space": "cosine", "product": group.product})
        else:
            collection = None

        for path in files:
            text = extract_text(path)
            digest = file_hash(path)
            chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)
            file_report = {"source": path.name, "path": str(path), "sha256": digest, "text_chars": len(text), "chunks": len(chunks)}
            product_report["files"].append(file_report)
            product_report["total_chunks"] += len(chunks)
            if dry_run:
                continue
            assert collection is not None
            assert client is not None
            for idx, chunk in enumerate(chunks):
                chunk_id = f"{group.product}:{path.name}:{digest[:12]}:{idx}"
                collection.add(
                    ids=[chunk_id],
                    embeddings=[embed_text(client, chunk, embedding_model)],
                    documents=[chunk],
                    metadatas=[
                        {
                            "product": group.product,
                            "source": path.name,
                            "source_path": str(path),
                            "chunk_index": idx,
                            "file_hash": digest,
                        }
                    ],
                )
        report["products"][group.product] = product_report
    return report


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Rebuild WESS chatbot ChromaDB from product documents.")
    parser.add_argument("--dry-run", action="store_true", help="Extract and count chunks without writing ChromaDB")
    parser.add_argument("--chunk-size", type=int, default=900)
    parser.add_argument("--overlap", type=int, default=120)
    parser.add_argument("--report", default="ingest_report.json", help="Path to write JSON report")
    parser.add_argument("--allow-partial", action="store_true", help="Allow destructive rebuild even when configured source folders are missing")
    args = parser.parse_args(argv)

    report = rebuild_database(dry_run=args.dry_run, allow_partial=args.allow_partial, chunk_size=args.chunk_size, overlap=args.overlap)
    report_path = Path(args.report)
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    for product, info in report["products"].items():
        print(f"{product}: {info['total_chunks']} chunks from {len(info['files'])} files")
    print(f"report: {report_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
